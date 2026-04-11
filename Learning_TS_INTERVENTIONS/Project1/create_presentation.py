"""
Generate HW5 Presentation as PPTX
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title
    left, top, width, height = Inches(0.5), Inches(2), Inches(9), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    left, top = Inches(0.5), Inches(3.5)
    txBox2 = slide.shapes.add_textbox(left, top, width, Inches(2))
    tf2 = txBox2.text_frame
    for line in subtitle.split('\n'):
        p = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
        p.text = line
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets, subbullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title
    left, top, width, height = Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    # Content
    left, top, width, height = Inches(0.4), Inches(1.2), Inches(9.2), Inches(5.5)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf2.add_paragraph() if i > 0 or tf2.paragraphs[0].text else tf2.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(18)
        p.level = 0
        p.space_after = Pt(6)

        if subbullets and i < len(subbullets) and subbullets[i]:
            for sub in subbullets[i]:
                p2 = tf2.add_paragraph()
                p2.text = "    • " + sub
                p2.font.size = Pt(16)
                p2.level = 1

def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    left, top, width, height = Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.7)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    left, top = Inches(0.3), Inches(1.1)
    width, height = Inches(9.4), Inches(0.4 * num_rows)
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(13)

    return slide

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
        "Time Series Prediction of Housing Prices",
        "Project 1 — Results & Takeaways\n\nDataset: Zillow House Prices (Boston, MA)\nTraining: 2010–2017 | Testing: 2018–2019\nExogenous Variable: 30-Year Mortgage Rate")

    # Slide 2: Assumptions
    slide = add_table_slide(prs, "1. Assumptions We Made",
        ["Assumption", "Why We Made It", "How We Tested It"],
        [
            ["Stationarity", "AR/ARMA requires stationary series", "ADF: raw p=0.98 → detrended p=0.03"],
            ["Linearity", "Model class is linear in lags", "Implicit in AR/ARX/ARMAX"],
            ["Credit↔Housing", "Mortgage rates affect prices", "Correlation: Δprice vs rate = −0.40***"],
            ["Finite AR order", "Parsimony; AIC/BIC penalize complexity", "Compared p=1,2,3,4,5 via AIC/BIC"],
        ])
    # Add note
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(4.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Key insight: We TESTED stationarity rather than assumed it."
    p.font.size = Pt(16)
    p.font.bold = True

    # Slide 3: Testing Credit-Housing Assumption
    slide = add_table_slide(prs, "1b. Testing the Credit-Housing Assumption",
        ["Relationship", "Correlation", "p-value", "Interpretation"],
        [
            ["Δ price vs rate level", "−0.40", "0.0000", "Higher rates → lower price growth"],
            ["Δ price vs Δ rate", "0.17", "0.11", "Change in rate less predictive"],
        ])
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(3.5), Inches(9), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    bullets = [
        "Regression: Δprice = 6741 − 1382 × rate  (R² = 0.16, p < 0.001)",
        "Each 1 p.p. rate increase → ~$1,382 less monthly price growth",
        "This motivated including interest rate as exogenous variable in ARX/ARMAX"
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(16)

    # Slide 4: Trend Removal
    add_content_slide(prs, "2. Trend Removal",
        [
            "What We Did: OLS Linear Detrend",
            "Result: Raw ADF p=0.98 (unit root) → Detrended ADF p=0.03 (stationary)",
            "Other Methods Available:",
            "Future Choice:"
        ],
        [
            ["Regress price on time index, use residuals as detrended series"],
            ["Successfully achieved stationarity for AR/ARMA modeling"],
            ["Differencing (yt − yt−1): good for stochastic trend / I(1)",
             "Polynomial detrend: for curved trends (we tried quadratic, cubic)",
             "HP filter / splines: smooth, flexible trend extraction"],
            ["Linear OLS for interpretability",
             "Differencing if ADF strongly suggests I(1)"]
        ])

    # Slide 5: Generalization
    add_content_slide(prs, "3. Generalizing to a New Dataset",
        [
            "Step 1: EXPLORE — Plot series, check missing values, define prediction task",
            "Step 2: STATIONARITY — Run ADF test; if non-stationary → detrend/difference",
            "Step 3: CORRELATIONS — Cross-correlation, ACF/PACF, test lag significance",
            "Step 4: MODEL SELECTION — Start AR → ARMA → ARX/ARMAX",
            "Step 5: ORDER SELECTION — AIC/BIC + out-of-sample MSE",
            "Step 6: EVALUATE — Hold-out test set; report MSE/RMSE/MAE"
        ])

    # Slide 6: Model Choice
    add_content_slide(prs, "3b. Model Choice & Domain Knowledge",
        [
            "Which Model?",
            "Does Dataset Type Matter? YES",
            "Does Domain Knowledge Matter? YES"
        ],
        [
            ["AR/ARMA: when only single time series available",
             "ARX/ARMAX: when plausible exogenous driver exists (e.g., mortgage rate)"],
            ["Frequency matters (daily vs monthly)",
             "Trend/seasonality structure",
             "Availability of aligned exogenous variables"],
            ["Choice of exogenous variable (mortgage rate for housing)",
             "Whether to model levels vs. growth rates",
             "Plausible forecast horizon (macro shocks limit predictability)"]
        ])

    # Slide 7: Evaluation Metrics
    slide = add_table_slide(prs, "4. Evaluation Metrics",
        ["Metric", "Purpose", "How We Used It"],
        [
            ["MSE", "Penalizes large errors", "Primary metric for comparison"],
            ["RMSE", "Same scale as variable", "Easier interpretation"],
            ["MAE", "Robust to outliers", "Alternative accuracy measure"],
            ["AIC/BIC", "Balance fit vs complexity", "AR order selection (optimal p=2 by BIC)"],
        ])
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(4.2), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "In-sample: AIC/BIC for order selection | Out-of-sample: MSE with rolling 1-step-ahead"
    p.font.size = Pt(16)
    p.font.bold = True

    # Slide 8: ARX vs AR
    slide = add_table_slide(prs, "5. ARX vs AR: Did Exogenous Inputs Help?",
        ["Model", "MSE Train (thousands)", "MSE Test (thousands)"],
        [
            ["AR(4)", "41,278", "282.9"],
            ["ARX(4)", "40,592", "281.7 ✓"],
        ])
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(3.2), Inches(9), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Result: ARX improved both in-sample and out-of-sample MSE (~0.4% improvement)"
    p.font.size = Pt(18)
    p.font.bold = True
    bullets = [
        "Why ARX helps: Interest rates capture credit channel affecting prices",
        "Why limited help: Relationship may already be in recent price lags",
        "Best ARMAX: (p=2, q=1) with test MSE = 300.6K"
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(16)

    # Slide 9: Error Sources
    slide = add_table_slide(prs, "6. Main Sources of Error",
        ["Source", "Description"],
        [
            ["Data", "Measurement error; interest rate resampled weekly→monthly"],
            ["Trend", "Linear trend may miss structural breaks"],
            ["Model", "AR assumes constant coefficients; possible regime changes"],
            ["Exogenous", "Interest rate may be endogenous; lag structure uncertain"],
            ["Evaluation", "Single train/test split; short test period (24 months)"],
        ])
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Mitigation: Cross-validation, multiple regions, longer test periods, regime-switching models"
    p.font.size = Pt(14)
    p.font.italic = True

    # Slide 10: Forecast Horizon
    slide = add_table_slide(prs, "7. How Far Ahead Can We Predict?",
        ["Strategy", "MSE Test (thousands)", "Interpretation"],
        [
            ["Short-term (rolling 1-step)", "294", "Reasonable accuracy"],
            ["Long-term (fixed multi-step)", "88,003", "300× worse — errors compound"],
        ])
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(3.2), Inches(9), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Comfortable predicting: Few months to ~1 year only"
    p.font.size = Pt(20)
    p.font.bold = True
    bullets = [
        "Error accumulation: Each step adds forecast error",
        "Exogenous shocks: Policy changes, crises unpredictable",
        "Parameter uncertainty: Extrapolation amplifies estimation error",
        "Recommendation: Re-fit regularly; use short-term forecasts; be cautious beyond 6-12 months"
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(16)

    # Slide 11: Summary
    slide = add_table_slide(prs, "Summary of Key Takeaways",
        ["Question", "Our Answer"],
        [
            ["Assumptions", "Tested stationarity (ADF), validated credit-price link (r=−0.40)"],
            ["Trend removal", "OLS linear detrend; differencing if strong I(1) evidence"],
            ["Generalization", "ADF → ACF/PACF → AR/ARX → AIC/BIC + out-of-sample MSE"],
            ["Metrics", "MSE for accuracy; AIC/BIC for model selection"],
            ["ARX vs AR", "ARX slightly better (~0.4%); include rate when available"],
            ["Error sources", "Data, trend, model assumptions, evaluation design"],
            ["Forecast horizon", "Comfortable only few months to ~1 year"],
        ])

    # Slide 12: Questions
    add_title_slide(prs, "Questions?",
        "Key Numbers:\n• Correlation (Δprice, rate): −0.40 (p < 0.001)\n• Short-term MSE: 294K vs Long-term: 88,003K (300× worse)\n• ARX(4) test MSE: 281.7K vs AR(4): 282.9K\n• ADF p-value: Raw 0.98 → Detrended 0.03")

    # Save
    prs.save('./HW5_Presentation.pptx')
    print("Presentation saved to HW5_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
